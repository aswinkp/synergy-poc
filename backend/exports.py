from __future__ import annotations

import csv
import re
import uuid
from pathlib import Path
from typing import Any

from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .config import EXPORTS_PATH

SUPPORTED_EXPORTS = {"csv", "xlsx", "pptx"}


def _columns(rows: list[dict[str, Any]]) -> list[str]:
    columns: list[str] = []
    for row in rows:
        for key in row:
            if key not in columns:
                columns.append(key)
    return columns


def _safe_cell(value: Any) -> Any:
    """Prevent spreadsheet software from treating source text as a formula."""
    if isinstance(value, str) and value.lstrip().startswith(("=", "+", "-", "@")):
        return f"'{value}"
    return value


def _safe_name(title: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9_-]+", "-", title.strip()).strip("-").lower()
    return cleaned[:60] or "synergy-export"


def _write_csv(path: Path, rows: list[dict[str, Any]], columns: list[str]) -> None:
    with path.open("w", encoding="utf-8-sig", newline="") as stream:
        writer = csv.DictWriter(stream, fieldnames=columns, extrasaction="ignore")
        if columns:
            writer.writeheader()
            for row in rows:
                writer.writerow({key: _safe_cell(row.get(key)) for key in columns})


def _write_xlsx(path: Path, rows: list[dict[str, Any]], columns: list[str], title: str) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = (_safe_name(title).replace("-", " ").title() or "Export")[:31]

    if columns:
        sheet.append(columns)
        header_fill = PatternFill("solid", fgColor="173245")
        for cell in sheet[1]:
            cell.fill = header_fill
            cell.font = Font(color="FFFFFF", bold=True)
            cell.alignment = Alignment(vertical="center")

        for row in rows:
            sheet.append([_safe_cell(row.get(key)) for key in columns])

        sheet.freeze_panes = "A2"
        sheet.auto_filter.ref = sheet.dimensions
        for index, column in enumerate(columns, start=1):
            longest = max(
                len(str(column)),
                *(len(str(_safe_cell(row.get(column)) or "")) for row in rows),
            )
            sheet.column_dimensions[get_column_letter(index)].width = min(max(longest + 2, 12), 42)

    workbook.save(path)
    workbook.close()


INK = RGBColor(14, 34, 51)
MUTED = RGBColor(91, 108, 120)
PANEL = RGBColor(241, 245, 247)
RULE = RGBColor(211, 221, 226)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(25, 165, 140)
BLUE = RGBColor(66, 153, 225)
SKY = RGBColor(116, 202, 230)
AMBER = RGBColor(244, 183, 64)
CORAL = RGBColor(235, 112, 91)
CHART_COLORS = (TEAL, BLUE, SKY, AMBER, CORAL, RGBColor(126, 105, 171))


def _add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: float,
    color: RGBColor = INK,
    bold: bool = False,
    font: str = "Aptos",
    align=PP_ALIGN.LEFT,
    vertical=MSO_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = vertical
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return shape


def _add_slide_number(slide, number: int) -> None:
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(7.08), Inches(11.95), Inches(0.012))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RULE
    rule.line.fill.background()
    _add_text(slide, "SYNERGY  /  LEARNING INTELLIGENCE", 0.72, 7.14, 4.8, 0.18, size=8.5, color=MUTED, bold=True)
    _add_text(slide, str(number), 12.0, 7.12, 0.62, 0.2, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def _add_slide_title(slide, title: str, number: int, eyebrow: str = "EXECUTIVE BRIEFING") -> None:
    _add_text(slide, eyebrow.upper(), 0.72, 0.38, 4.8, 0.22, size=10, color=TEAL, bold=True)
    single_line_title = " ".join(title.split())
    title_size = 35 if len(single_line_title) <= 46 else 32
    _add_text(slide, single_line_title, 0.72, 0.72, 11.9, 0.62, size=title_size, color=INK, bold=True, font="Aptos Display")
    _add_slide_number(slide, number)


def _summary_blocks(summary: str) -> list[str]:
    cleaned = [re.sub(r"^[#>*\-\s]+", "", part).strip() for part in summary.splitlines()]
    blocks = [part for part in cleaned if part]
    if not blocks and summary.strip():
        blocks = [summary.strip()]

    split_blocks: list[str] = []
    for block in blocks:
        if len(block) <= 520:
            split_blocks.append(block)
            continue
        sentences = re.split(r"(?<=[.!?])\s+", block)
        current = ""
        for sentence in sentences:
            candidate = f"{current} {sentence}".strip()
            if current and len(candidate) > 480:
                split_blocks.append(current)
                current = sentence
            else:
                current = candidate
        if current:
            split_blocks.append(current)
    return split_blocks or ["The requested analysis is ready for management review."]


def _summary_pages(summary: str) -> list[list[str]]:
    pages: list[list[str]] = []
    current: list[str] = []
    current_characters = 0
    for block in _summary_blocks(summary):
        if current and (len(current) >= 5 or current_characters + len(block) > 1_450):
            pages.append(current)
            current = []
            current_characters = 0
        current.append(block)
        current_characters += len(block)
    if current:
        pages.append(current)
    return pages


def _insight_excerpt(text: str, limit: int = 96) -> str:
    first_sentence = re.split(r"(?<=[.!?])\s+", text.strip(), maxsplit=1)[0]
    if len(first_sentence) <= limit:
        return first_sentence
    colon_clause = first_sentence.split(":", 1)[0].strip()
    if 35 <= len(colon_clause) <= limit:
        return colon_clause
    comma_clause = first_sentence.split(",", 1)[0].strip()
    if 35 <= len(comma_clause) <= limit:
        return comma_clause
    shortened = first_sentence[: limit + 1].rsplit(" ", 1)[0].rstrip(" ,;:")
    return f"{shortened}…"


def _add_briefing_slide(presentation, blocks: list[str], number: int, page_index: int) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = "Executive briefing" if page_index == 0 else "Evidence and implications"
    _add_slide_title(slide, title, number)

    lead_height = 1.62 if len(blocks) > 1 else 3.5
    lead_panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.72),
        Inches(1.62),
        Inches(11.9),
        Inches(lead_height),
    )
    lead_panel.fill.solid()
    lead_panel.fill.fore_color.rgb = PANEL
    lead_panel.line.fill.background()
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.72),
        Inches(1.62),
        Inches(0.11),
        Inches(lead_height),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    _add_text(
        slide,
        blocks[0],
        1.08,
        1.9,
        11.0,
        lead_height - 0.52,
        size=26 if len(blocks[0]) < 240 else 21,
        color=INK,
        bold=True,
        vertical=MSO_ANCHOR.MIDDLE,
    )

    remaining = blocks[1:]
    if not remaining:
        _add_text(
            slide,
            "Prepared from the connected workforce and learning data",
            1.08,
            5.56,
            8.5,
            0.35,
            size=16,
            color=MUTED,
        )
        return

    available_height = 3.18
    row_height = available_height / len(remaining)
    for index, block in enumerate(remaining):
        y = 3.58 + index * row_height
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), Inches(y + 0.07), Inches(0.16), Inches(0.16))
        marker.fill.solid()
        marker.fill.fore_color.rgb = CHART_COLORS[index % len(CHART_COLORS)]
        marker.line.fill.background()
        _add_text(
            slide,
            block,
            1.08,
            y,
            11.35,
            row_height - 0.08,
            size=17 if len(block) < 360 else 16,
            color=INK,
            vertical=MSO_ANCHOR.MIDDLE,
        )


def _chart_kind(chart_type: str, chart_rows: list[dict[str, Any]], label_key: str):
    if chart_type == "bar":
        labels = [str(row.get(label_key, "")) for row in chart_rows]
        return XL_CHART_TYPE.BAR_CLUSTERED if len(labels) > 6 or any(len(label) > 16 for label in labels) else XL_CHART_TYPE.COLUMN_CLUSTERED
    return {
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "area": XL_CHART_TYPE.AREA,
        "pie": XL_CHART_TYPE.DOUGHNUT,
    }[chart_type]


def _style_chart(chart, chart_type: str, value_keys: list[str]) -> None:
    chart.has_title = False
    chart.has_legend = len(value_keys) > 1 or chart_type == "pie"
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = "Aptos"
        chart.legend.font.size = Pt(11)

    for series_index, series in enumerate(chart.series):
        color = CHART_COLORS[series_index % len(CHART_COLORS)]
        if chart_type == "line":
            series.format.line.color.rgb = color
            series.format.line.width = Pt(2.5)
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
            series.format.line.fill.background()

    plot = chart.plots[0]
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.font.name = "Aptos"
    labels.font.size = Pt(10)
    labels.font.color.rgb = INK
    labels.show_legend_key = False
    if chart_type == "pie":
        plot.hole_size = 58
        plot.vary_by_categories = True
        labels.show_category_name = True
        labels.show_percentage = True
        labels.show_value = False
        labels.position = XL_LABEL_POSITION.BEST_FIT
        if chart.series:
            for point_index, point in enumerate(chart.series[0].points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = CHART_COLORS[point_index % len(CHART_COLORS)]
                point.format.line.fill.background()
        return

    labels.show_value = True
    labels.position = XL_LABEL_POSITION.OUTSIDE_END
    category_axis = chart.category_axis
    category_axis.tick_labels.font.name = "Aptos"
    category_axis.tick_labels.font.size = Pt(10)
    category_axis.format.line.color.rgb = RULE
    value_axis = chart.value_axis
    value_axis.tick_labels.font.name = "Aptos"
    value_axis.tick_labels.font.size = Pt(10)
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = RULE
    value_axis.format.line.fill.background()


def _add_chart_slide(
    presentation,
    visualization: dict[str, Any],
    rows: list[dict[str, Any]],
    summary: str,
    number: int,
) -> None:
    chart_type = str(visualization["type"])
    chart_rows = visualization.get("data", rows)
    value_keys = list(visualization.get("valueKeys", []))
    label_key = str(visualization.get("labelKey", "label"))
    chart_data = CategoryChartData()
    chart_data.categories = [str(row.get(label_key, "")) for row in chart_rows]
    for key in value_keys:
        chart_data.add_series(
            key.replace("_", " ").title(),
            [float(row.get(key) or 0) for row in chart_rows],
        )

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_slide_title(slide, str(visualization.get("title") or "Key evidence"), number, "DATA STORY")
    chart = slide.shapes.add_chart(
        _chart_kind(chart_type, chart_rows, label_key),
        Inches(0.68),
        Inches(1.55),
        Inches(8.35),
        Inches(5.15),
        chart_data,
    ).chart
    _style_chart(chart, chart_type, value_keys)

    insight_panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(9.35),
        Inches(1.55),
        Inches(3.28),
        Inches(5.15),
    )
    insight_panel.fill.solid()
    insight_panel.fill.fore_color.rgb = PANEL
    insight_panel.line.fill.background()
    _add_text(slide, "MANAGEMENT TAKEAWAYS", 9.7, 1.92, 2.6, 0.3, size=11, color=TEAL, bold=True)
    insights = _summary_blocks(summary)[:3]
    for index, insight in enumerate(insights):
        y = 2.55 + index * 1.28
        _add_text(slide, f"0{index + 1}", 9.7, y, 0.42, 0.22, size=10, color=CHART_COLORS[index], bold=True)
        _add_text(
            slide,
            _insight_excerpt(insight),
            9.7,
            y + 0.28,
            2.55,
            0.82,
            size=13,
            color=INK,
        )


def _add_data_slide(
    presentation,
    rows: list[dict[str, Any]],
    columns: list[str],
    number: int,
) -> None:
    supporting_rows = rows[:7]
    supporting_columns = columns[:6]
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_slide_title(slide, "Evidence behind the recommendation", number, "SUPPORTING DATA")
    _add_text(
        slide,
        "A concise view of the data used in the analysis",
        0.72,
        1.38,
        8.5,
        0.3,
        size=15,
        color=MUTED,
    )
    table_shape = slide.shapes.add_table(
        len(supporting_rows) + 1,
        len(supporting_columns),
        Inches(0.72),
        Inches(1.86),
        Inches(11.9),
        Inches(4.72),
    )
    table = table_shape.table
    first_width = 3.05 if len(supporting_columns) > 1 else 11.9
    table.columns[0].width = Inches(first_width)
    if len(supporting_columns) > 1:
        remaining_width = (11.9 - first_width) / (len(supporting_columns) - 1)
        for column_index in range(1, len(supporting_columns)):
            table.columns[column_index].width = Inches(remaining_width)

    for column_index, column in enumerate(supporting_columns):
        cell = table.cell(0, column_index)
        cell.text = column.replace("_", " ").title()
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = "Aptos"
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.size = Pt(12)

    for row_index, row in enumerate(supporting_rows, start=1):
        for column_index, column in enumerate(supporting_columns):
            cell = table.cell(row_index, column_index)
            value = _safe_cell(row.get(column))
            cell.text = "" if value is None else str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 else PANEL
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Aptos"
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = INK

    if len(rows) > len(supporting_rows) or len(columns) > len(supporting_columns):
        _add_text(
            slide,
            f"Showing {len(supporting_rows):,} of {len(rows):,} rows and {len(supporting_columns):,} of {len(columns):,} columns. "
            "Request CSV or Excel for the complete dataset.",
            0.74,
            6.68,
            11.2,
            0.25,
            size=10,
            color=MUTED,
        )


def _write_pptx(
    path: Path,
    rows: list[dict[str, Any]],
    columns: list[str],
    title: str,
    summary: str,
    visualization: dict[str, Any] | None,
) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, presentation.slide_width, presentation.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = INK
    background.line.fill.background()
    color_field = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(9.78),
        Inches(0),
        Inches(3.56),
        Inches(7.5),
    )
    color_field.fill.solid()
    color_field.fill.fore_color.rgb = TEAL
    color_field.line.fill.background()
    _add_text(title_slide, "SYNERGY", 0.78, 0.58, 2.3, 0.35, size=12, color=WHITE, bold=True)
    _add_text(title_slide, "LEARNING INTELLIGENCE", 10.15, 0.58, 2.45, 0.35, size=10, color=WHITE, bold=True)
    _add_text(
        title_slide,
        title,
        0.78,
        2.02,
        8.35,
        3.25,
        size=50 if len(title) < 72 else 42,
        color=WHITE,
        bold=True,
        font="Aptos Display",
        vertical=MSO_ANCHOR.MIDDLE,
    )
    _add_text(
        title_slide,
        "A management briefing generated from the connected workforce and learning data",
        0.82,
        6.28,
        7.9,
        0.55,
        size=16,
        color=RGBColor(194, 207, 216),
    )

    slide_number = 2
    for page_index, blocks in enumerate(_summary_pages(summary)):
        _add_briefing_slide(presentation, blocks, slide_number, page_index)
        slide_number += 1

    chart_type = visualization.get("type") if visualization else None
    chart_rows = visualization.get("data", rows) if visualization else rows
    value_keys = visualization.get("valueKeys", []) if visualization else []
    supported_charts = {"bar", "line", "area", "pie"}
    if chart_type in supported_charts and chart_rows and value_keys:
        _add_chart_slide(presentation, visualization, rows, summary, slide_number)
        slide_number += 1

    if rows and columns:
        _add_data_slide(presentation, rows, columns, slide_number)

    presentation.save(path)


def create_export(
    rows: list[dict[str, Any]],
    title: str,
    export_format: str,
    *,
    summary: str = "",
    visualization: dict[str, Any] | None = None,
) -> dict[str, Any]:
    if export_format not in SUPPORTED_EXPORTS:
        raise ValueError(f"Unsupported export format: {export_format}")

    EXPORTS_PATH.mkdir(parents=True, exist_ok=True)
    export_id = str(uuid.uuid4())
    path = EXPORTS_PATH / f"{export_id}.{export_format}"
    columns = _columns(rows)
    if export_format == "csv":
        _write_csv(path, rows, columns)
    elif export_format == "xlsx":
        _write_xlsx(path, rows, columns, title)
    else:
        _write_pptx(path, rows, columns, title, summary, visualization)

    display_name = f"{_safe_name(title)}.{export_format}"
    return {
        "id": export_id,
        "format": export_format,
        "filename": display_name,
        "url": f"/api/exports/{export_id}",
        "row_count": len(rows),
        "size_bytes": path.stat().st_size,
        "title": title,
    }


def attachment_path(attachment: dict[str, Any]) -> Path | None:
    try:
        export_id = str(uuid.UUID(str(attachment["id"])))
        export_format = str(attachment["format"])
    except (KeyError, TypeError, ValueError):
        return None
    if export_format not in SUPPORTED_EXPORTS:
        return None
    return EXPORTS_PATH / f"{export_id}.{export_format}"
