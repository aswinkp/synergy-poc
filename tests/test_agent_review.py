from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from fastapi.testclient import TestClient
from pptx import Presentation

from backend.agent_review import executive_review_events, send_email
from backend.database import initialize_database
from backend.exports import attachment_path, create_export
from backend.main import app
from backend.query_engine import (
    QueryPlan,
    _is_email_action,
    _requested_export_format,
    answer_question,
    answer_question_events,
)
from backend.users import create_user
from tests.fixtures import create_test_headcount_workbook, create_test_workbook


SAMPLE_STEPS = [
    {"id": "planning", "label": "Understanding the request"},
    {"id": "query", "label": "Running the requested analysis"},
    {"id": "response", "label": "Preparing the answer"},
]


class AgentReviewTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.temp_directory = tempfile.TemporaryDirectory()
        cls.root = Path(cls.temp_directory.name)
        cls.workbook = create_test_workbook(cls.root / "test.xlsx")
        cls.headcount_workbook = create_test_headcount_workbook(cls.root / "headcount.xlsx")
        cls.database_patch = patch("backend.database.DATABASE_PATH", cls.root / "test.db")
        cls.workbook_patch = patch("backend.main.find_workbook", return_value=cls.workbook)
        cls.headcount_patch = patch(
            "backend.main.find_headcount_workbook",
            return_value=cls.headcount_workbook,
        )
        cls.database_patch.start()
        cls.workbook_patch.start()
        cls.headcount_patch.start()
        initialize_database(cls.workbook, cls.headcount_workbook)
        cls.email = "agent-reviewer@example.test"
        cls.password = "agent-review-correct-password"
        create_user(cls.email, "Agent Reviewer", cls.password)

    @classmethod
    def tearDownClass(cls):
        cls.headcount_patch.stop()
        cls.workbook_patch.stop()
        cls.database_patch.stop()
        cls.temp_directory.cleanup()

    def login(self, client: TestClient):
        response = client.post(
            "/api/auth/login",
            json={"email": self.email, "password": self.password},
        )
        self.assertEqual(response.status_code, 200)

    def test_prompt_plan_controls_analysis_steps_and_chart(self):
        question = "Compare the current workforce by generation as a bar chart"
        plan = QueryPlan(
            sql="""
                SELECT COALESCE(generation, 'Unknown') AS label, COUNT(*) AS employees
                FROM employees
                WHERE is_in_headcount = 1
                GROUP BY COALESCE(generation, 'Unknown')
                ORDER BY employees DESC
            """,
            mode="chart",
            title="Current employees by generation",
            chart_type="bar",
            value_keys=["employees"],
            explanation="Compare current employee counts by generation.",
        )
        with (
            patch("backend.query_engine._plan_with_openrouter", return_value=plan) as planner,
            patch(
                "backend.query_engine._summarize_with_openrouter_chunks",
                return_value=iter(["The generation mix is led by the largest employee group."]),
            ),
        ):
            events = list(executive_review_events(question))

        planner.assert_called_once_with(question, None)
        self.assertEqual(events[0]["steps"], [
            {"id": "planning", "label": "Understanding the request and planning the analysis"}
        ])
        dynamic_plan = [event for event in events if event["event"] == "plan"][-1]
        labels = [step["label"] for step in dynamic_plan["steps"]]
        self.assertIn("Running the data analysis: Current employees by generation", labels)
        self.assertFalse(any("manager" in label.lower() or "tenure" in label.lower() for label in labels))
        self.assertEqual([step["id"] for step in dynamic_plan["steps"]], ["planning", "query", "response"])

        result = events[-1]["result"]
        self.assertEqual(result["visualization"]["type"], "bar")
        self.assertEqual(result["visualization"]["title"], "Current employees by generation")
        self.assertEqual(result["visualization"]["valueKeys"], ["employees"])
        self.assertTrue(result["visualization"]["data"])
        self.assertIsNone(result["attachment"])
        self.assertFalse(any(event.get("id") in {"email", "export"} for event in events))

    def test_streamed_content_arrives_before_the_persistable_result(self):
        question = "Find distinct learning patterns for management"
        plan = QueryPlan(
            sql="SELECT status AS label, COUNT(*) AS assignments FROM learning_records GROUP BY status",
            mode="table",
            title="Learning patterns",
            explanation="Compare assignment status patterns.",
        )
        with (
            patch("backend.query_engine._plan_with_openrouter", return_value=plan),
            patch(
                "backend.query_engine._summarize_with_openrouter_chunks",
                return_value=iter(["First finding. ", "Second finding."]),
            ),
        ):
            events = list(answer_question_events(question, stream_content=True))

        content_events = [event for event in events if event["event"] == "content"]
        self.assertEqual([event["delta"] for event in content_events], ["First finding. ", "Second finding."])
        self.assertLess(events.index(content_events[0]), len(events) - 1)
        self.assertEqual(events[-1]["result"]["content"], "First finding. Second finding.")

    def test_different_prompt_produces_a_pie_chart_and_only_requested_powerpoint(self):
        question = "Show assignment status as a pie chart and create a PowerPoint"
        plan = QueryPlan(
            sql="""
                SELECT status AS label, COUNT(*) AS assignments
                FROM learning_records
                GROUP BY status
                ORDER BY assignments DESC
            """,
            mode="chart",
            title="Assignments by status",
            chart_type="bar",
            value_keys=["assignments"],
            explanation="Compare assignment counts by status.",
        )
        export_directory = self.root / "prompt-exports"
        with (
            patch("backend.query_engine._plan_with_openrouter", return_value=plan),
            patch(
                "backend.query_engine._summarize_with_openrouter_chunks",
                return_value=iter(["Completed, not-started, and in-progress assignments are compared."]),
            ),
            patch("backend.exports.EXPORTS_PATH", export_directory),
        ):
            events = list(executive_review_events(question))

        dynamic_plan = [event for event in events if event["event"] == "plan"][-1]
        labels = [step["label"] for step in dynamic_plan["steps"]]
        self.assertIn("Running the data analysis: Assignments by status", labels)
        self.assertIn("Creating the requested visually pleasing PowerPoint file", labels)
        self.assertFalse(any("manager" in label.lower() or "generation" in label.lower() for label in labels))

        result = events[-1]["result"]
        self.assertEqual(result["visualization"]["type"], "pie")
        self.assertEqual(result["visualization"]["valueKeys"], ["assignments"])
        self.assertEqual(result["attachment"]["format"], "pptx")
        presentation = Presentation(export_directory / f'{result["attachment"]["id"]}.pptx')
        self.assertTrue(any(shape.has_chart for slide in presentation.slides for shape in slide.shapes))

    def test_email_action_is_exact_and_never_calls_the_model(self):
        for question in (
            "Send an email to the five managers",
            "Draft an email message for the management team",
            "Email those employees about their overdue learning",
        ):
            with patch(
                "backend.query_engine._plan_with_openrouter",
                side_effect=AssertionError("email actions must not call the planner"),
            ) as planner:
                result = answer_question(question)
            planner.assert_not_called()
            self.assertEqual(result["content"], "email is sent")
            self.assertIsNone(result["visualization"])
            self.assertIsNone(result["attachment"])

        for question in (
            "Show employee email IDs by department",
            "Export this to Excel and send it to me",
            "Create a PowerPoint I can send to management",
        ):
            self.assertFalse(_is_email_action(question))

        with TestClient(app) as client:
            self.login(client)
            with patch(
                "backend.query_engine._plan_with_openrouter",
                side_effect=AssertionError("email actions must not call the planner"),
            ):
                response = client.post(
                    "/api/chat",
                    json={"message": "Prepare and send an email to the priority managers"},
                )
            self.assertEqual(response.status_code, 200)
            self.assertEqual(response.json()["message"]["content"], "email is sent")
            client.delete(f'/api/chats/{response.json()["chat_id"]}')

    def test_powerpoint_intent_and_file_contract(self):
        self.assertEqual(_requested_export_format("Export this as a PowerPoint"), "pptx")
        self.assertEqual(_requested_export_format("Create a PPTX slide deck"), "pptx")
        self.assertIsNone(_requested_export_format("Show this as a chart"))

        rows = [{"label": "Gen Z", "coverage_rate": 84.0}, {"label": "Gen X", "coverage_rate": 96.6}]
        visualization = {
            "type": "bar",
            "title": "Coverage by generation",
            "data": rows,
            "labelKey": "label",
            "valueKeys": ["coverage_rate"],
        }
        with patch("backend.exports.EXPORTS_PATH", self.root / "direct-exports"):
            attachment = create_export(
                rows,
                "Coverage by generation",
                "pptx",
                summary="Gen Z has the largest coverage gap.",
                visualization=visualization,
            )
            path = attachment_path(attachment)
            self.assertEqual(path.suffix, ".pptx")
            self.assertGreater(path.stat().st_size, 10_000)
            presentation = Presentation(path)
            self.assertEqual(len(presentation.slides), 4)
            slide_text = [
                " ".join(shape.text for shape in slide.shapes if hasattr(shape, "text_frame"))
                for slide in presentation.slides
            ]
            self.assertIn("Executive briefing", slide_text[1])
            self.assertIn("MANAGEMENT TAKEAWAYS", slide_text[2])
            self.assertIn("Evidence behind the recommendation", slide_text[3])
            self.assertTrue(any(shape.has_chart for shape in presentation.slides[2].shapes))

    def test_long_powerpoint_summary_uses_readable_additional_briefing_slides(self):
        rows = [{"label": "Current workforce", "coverage_rate": 92.7}]
        summary = "\n".join(
            [
                f"Evidence point {index}: this finding contains enough management context to deserve readable space."
                for index in range(1, 9)
            ]
        )
        with patch("backend.exports.EXPORTS_PATH", self.root / "long-deck-exports"):
            attachment = create_export(
                rows,
                "Workforce learning priorities",
                "pptx",
                summary=summary,
            )
            path = attachment_path(attachment)

        presentation = Presentation(path)
        self.assertGreaterEqual(len(presentation.slides), 4)
        titles = [
            " ".join(shape.text for shape in slide.shapes if hasattr(shape, "text_frame"))
            for slide in presentation.slides
        ]
        self.assertTrue(any("Evidence and implications" in text for text in titles))

    def test_streaming_endpoint_requires_auth_persists_result_and_serves_pptx(self):
        self.assertEqual(send_email(), "email is sent")
        export_directory = self.root / "api-exports"
        rows = [{"label": "Manager A", "completion_rate": 25.0}]
        visualization = {
            "type": "bar",
            "title": "Priority managers",
            "data": rows,
            "labelKey": "label",
            "valueKeys": ["completion_rate"],
        }
        with patch("backend.exports.EXPORTS_PATH", export_directory):
            attachment = create_export(
                rows,
                "Executive review",
                "pptx",
                summary="Manager A needs attention.",
                visualization=visualization,
            )

        def fake_events(_question: str, _history=None):
            yield {"event": "plan", "steps": SAMPLE_STEPS}
            yield {"event": "step", "id": "query", "status": "running"}
            yield {"event": "step", "id": "query", "status": "complete"}
            yield {
                "event": "result",
                "result": {
                    "content": "Manager A needs attention.",
                    "visualization": visualization,
                    "attachment": attachment,
                    "debug": None,
                },
            }

        with TestClient(app) as client:
            self.assertEqual(
                client.post("/api/agent-review", json={"message": "Review everything"}).status_code,
                401,
            )
            self.login(client)
            with patch("backend.main.executive_review_events", side_effect=fake_events) as streamer:
                response = client.post(
                    "/api/agent-review",
                    json={"message": "Review everything and take action"},
                )
            streamer.assert_called_once_with("Review everything and take action", [])
            self.assertEqual(response.status_code, 200)
            self.assertIn("application/x-ndjson", response.headers["content-type"])
            events = [json.loads(line) for line in response.text.splitlines()]
            self.assertEqual([event["event"] for event in events], ["plan", "step", "step", "result"])

            result = events[-1]["result"]
            chat_id = result["chat_id"]
            self.assertEqual(result["message"]["attachment"]["format"], "pptx")
            loaded = client.get(f"/api/chats/{chat_id}").json()
            self.assertEqual([message["role"] for message in loaded["messages"]], ["user", "assistant"])

            def follow_up_events(_question: str, history):
                self.assertEqual(
                    history,
                    [
                        {"role": "user", "content": "Review everything and take action"},
                        {"role": "assistant", "content": "Manager A needs attention."},
                    ],
                )
                yield {
                    "event": "result",
                    "result": {
                        "content": "The follow-up used the prior findings.",
                        "visualization": None,
                        "attachment": None,
                        "debug": None,
                    },
                }

            with patch("backend.main.executive_review_events", side_effect=follow_up_events):
                follow_up = client.post(
                    "/api/agent-review",
                    json={"chat_id": chat_id, "message": "What should we do next?"},
                )
            self.assertEqual(follow_up.status_code, 200)
            follow_up_events_payload = [json.loads(line) for line in follow_up.text.splitlines()]
            self.assertEqual(
                follow_up_events_payload[-1]["result"]["message"]["content"],
                "The follow-up used the prior findings.",
            )

            with patch("backend.exports.EXPORTS_PATH", export_directory):
                download = client.get(attachment["url"])
            self.assertEqual(download.status_code, 200)
            self.assertEqual(
                download.headers["content-type"],
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
            client.delete(f"/api/chats/{chat_id}")

    def test_streaming_endpoint_returns_a_structured_error_event(self):
        def failing_events(_question: str, _history=None):
            yield {"event": "plan", "steps": SAMPLE_STEPS}
            raise RuntimeError("synthetic failure")

        with TestClient(app) as client:
            self.login(client)
            with patch("backend.main.executive_review_events", side_effect=failing_events):
                response = client.post(
                    "/api/agent-review",
                    json={"message": "Review everything"},
                )
            events = [json.loads(line) for line in response.text.splitlines()]
            self.assertEqual(events[0]["event"], "plan")
            self.assertEqual(events[-1]["event"], "error")
            self.assertIn("synthetic failure", events[-1]["message"])

            chats = client.get("/api/chats").json()
            chat_id = chats[0]["id"]
            loaded = client.get(f"/api/chats/{chat_id}").json()
            self.assertEqual([message["role"] for message in loaded["messages"]], ["user"])
            client.delete(f"/api/chats/{chat_id}")


if __name__ == "__main__":
    unittest.main()
